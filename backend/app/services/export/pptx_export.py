"""Render the slide editor's deck into .pptx."""
from __future__ import annotations

import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

ALIGNMENTS = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def _rgb(value: str | None) -> RGBColor | None:
    if not value:
        return None
    digits = value.strip().lstrip("#")
    if len(digits) == 3:
        digits = "".join(ch * 2 for ch in digits)
    if len(digits) != 6:
        return None
    try:
        return RGBColor.from_string(digits.upper())
    except ValueError:
        return None


def _style_frame(frame: Any, style: dict[str, Any] | None, default_size: int) -> None:
    style = style or {}
    color = _rgb(style.get("color"))
    size = Pt(float(style.get("fontSize") or default_size))
    align = ALIGNMENTS.get(style.get("align"))
    for paragraph in frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = size
            run.font.bold = bool(style.get("bold", run.font.bold))
            run.font.italic = bool(style.get("italic", run.font.italic))
            if style.get("underline"):
                run.font.underline = True
            if style.get("fontFamily"):
                run.font.name = style["fontFamily"]
            if color is not None:
                run.font.color.rgb = color


def _fill_bullets(placeholder: Any, slide: dict[str, Any]) -> None:
    frame = placeholder.text_frame
    frame.clear()
    frame.word_wrap = True

    lines: list[tuple[str, int]] = [(str(b), 0) for b in slide.get("bullets") or [] if str(b).strip()]
    body = str(slide.get("body") or "").strip()
    if body:
        lines.extend((line.strip(), 0) for line in body.splitlines() if line.strip())

    if not lines:
        frame.text = ""
        return

    first_text, _ = lines[0]
    frame.paragraphs[0].text = first_text
    for text, level in lines[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level


def deck_to_pptx(deck: dict[str, Any], title: str = "") -> bytes:
    presentation = Presentation()
    presentation.slide_width = Pt(960)   # 16:9
    presentation.slide_height = Pt(540)

    title_layout = presentation.slide_layouts[0]
    content_layout = presentation.slide_layouts[1]
    section_layout = presentation.slide_layouts[5]

    slides = deck.get("slides") or []
    if not slides:
        slides = [{"layout": "title", "title": title or deck.get("title") or "Presentation"}]

    for index, slide_data in enumerate(slides):
        layout_name = slide_data.get("layout") or ("title" if index == 0 else "bullets")
        has_content = bool(slide_data.get("bullets")) or bool(str(slide_data.get("body") or "").strip())

        if layout_name == "title":
            slide = presentation.slides.add_slide(title_layout)
        elif not has_content:
            slide = presentation.slides.add_slide(section_layout)
        else:
            slide = presentation.slides.add_slide(content_layout)

        heading = slide_data.get("title") or (title if index == 0 else "")
        if slide.shapes.title is not None:
            slide.shapes.title.text = str(heading)
            _style_frame(
                slide.shapes.title.text_frame,
                slide_data.get("titleStyle"),
                40 if layout_name == "title" else 30,
            )

        placeholders = [
            shape
            for shape in slide.placeholders
            if shape.placeholder_format.idx != 0 and shape.has_text_frame
        ]
        if placeholders:
            body_placeholder = placeholders[0]
            if layout_name == "title":
                body_placeholder.text_frame.text = str(
                    slide_data.get("subtitle") or slide_data.get("body") or ""
                )
                _style_frame(body_placeholder.text_frame, slide_data.get("bodyStyle"), 20)
            elif has_content:
                _fill_bullets(body_placeholder, slide_data)
                _style_frame(body_placeholder.text_frame, slide_data.get("bodyStyle"), 18)
            else:
                body_placeholder.text_frame.text = ""

        notes = str(slide_data.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()
